import os
import logging
import shutil
from pathlib import Path
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import ChatPromptTemplate
from langchain.docstore.document import Document as LangchainDocument

from docx import Document as DocxDocument
from pptx import Presentation
from fpdf import FPDF
# Hàm generate_text
def generate_text(input: str) -> str:
    return input + " generated"

# Cấu hình logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Cấu hình API key cho Google Generative AI
os.environ["GOOGLE_API_KEY"] = "AIzaSyC437JOnEhmdcIqtuHSJAEkzaD6faXwy-I"

# Khởi tạo mô hình Gemini
llm = ChatGoogleGenerativeAI(
    model="gemini-1.5-pro",
    temperature=0.5,
    max_tokens=4096,
    timeout=120,  # Tăng thời gian chờ
    max_retries=10,  # Tăng số lần thử lại
)


# Prompt tạo câu hỏi tự luận
essay_question_prompt = PromptTemplate(
    input_variables=["content", "num_questions"],
    template=(
        "Generate {num_questions} essay-style questions based on the following content:\n"
        "\n{content}\n\nFormat: Question text and sample answer."
    )
)

# Prompt tóm tắt nội dung
summarization_prompt = ChatPromptTemplate.from_messages(
    [
        ("system", "You are an expert summarizer. Summarize the following text."),
        ("human", "{content}"),
    ]
)

# Hàm đọc file .docx
def load_docx(file_path):
    doc = DocxDocument(file_path)
    content = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
    if not content.strip():
        raise ValueError(f"File {file_path} is empty or invalid.")
    return [LangchainDocument(page_content=content)]

# Hàm đọc file .pptx
def load_pptx(file_path):
    presentation = Presentation(file_path)
    content = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                content.append(shape.text)
    if not content:
        raise ValueError(f"File {file_path} is empty or invalid.")
    return [LangchainDocument(page_content="\n".join(content))]

# Loader cho các loại file
FILE_LOADERS = {
    ".pdf": PyPDFLoader,
    ".docx": load_docx,
    ".pptx": load_pptx,
}

# Hàm trích xuất nội dung từ file
def extract_text_from_file(file_path):
    file_extension = Path(file_path).suffix.lower()
    if file_extension not in FILE_LOADERS:
        raise ValueError(f"Unsupported file format: {file_extension}")
    
    loader = FILE_LOADERS[file_extension]
    documents = loader(file_path).load() if file_extension == ".pdf" else loader(file_path)
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    return text_splitter.split_documents(documents)
# Hàm tạo câu hỏi tự luận
def generate_essay_questions(content, num_questions):
    summarization_chain = summarization_prompt | llm
    summarized_content = summarization_chain.invoke({"content": content})
    
    summarized_text = (
        summarized_content.content if hasattr(summarized_content, "content") else str(summarized_content)
    ).strip()

    if not summarized_text:
        raise ValueError("Summarized content is empty. Please check the input content.")

    question_chain = essay_question_prompt | llm
    questions = question_chain.invoke(
        {
            "content": summarized_text,
            "num_questions": num_questions,
        }
    )

    questions_content = (
        questions.content if hasattr(questions, "content") else str(questions)
    ).strip()

    if not questions_content:
        raise ValueError("Generated questions are empty. Please check the LLM chain.")

    print("LLM Raw Output:\n", questions_content)  # Debug log

    return questions_content

# Hàm lưu định dạng PDF
def save_as_pdf(output, file_name):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", style="B", size=16)
    pdf.set_text_color(0, 0, 128)
    pdf.cell(0, 10, "Essay Questions", ln=True, align="C")
    pdf.ln(10)

    questions = output.split("**Question")
    for idx, question in enumerate(questions[1:], 1):
        question = question.strip()

        text_start = question.find("**Text:**")
        answer_start = question.find("**Sample Answer:**")

        if text_start != -1 and answer_start != -1:
            question_text = question[text_start + len("**Text:**"):answer_start].strip()
            sample_answer = question[answer_start + len("**Sample Answer:**"):].strip()

            pdf.set_text_color(0, 0, 0)
            pdf.set_font("Arial", style="B", size=12)
            pdf.cell(0, 10, f"Question {idx}:", ln=True)
            pdf.set_font("Arial", style="", size=12)
            pdf.multi_cell(0, 10, question_text)

            pdf.set_text_color(0, 0, 128)
            pdf.set_font("Arial", style="B", size=12)
            pdf.cell(0, 10, "Sample Answer:", ln=True)
            pdf.set_text_color(0, 0, 0)
            pdf.set_font("Arial", style="", size=12)
            pdf.multi_cell(0, 10, sample_answer)
            pdf.ln(10)
        else:
            pdf.set_text_color(255, 0, 0)
            pdf.cell(0, 10, f"Invalid Question Format for Question {idx}.", ln=True)

    pdf.output(file_name)
    print(f"PDF file saved as {file_name}.")

# Hàm lưu kết quả
def save_output(output, output_format, file_name):
    if isinstance(output, str):
        if output_format == "pdf":
            save_as_pdf(output, file_name)
        else:
            raise ValueError("Unsupported output format.")
    else:
        raise TypeError("Output must be a string.")

# Main workflow
def main():
    file_path = "./testing.docx"
    num_questions = 5
    output_file = "essay_questions.pdf"

    try:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File does not exist: {file_path}")

        extracted_content = extract_text_from_file(file_path)
        full_content = "\n".join([doc.page_content for doc in extracted_content])

        generated_questions = generate_essay_questions(full_content, num_questions)
        
        # Sử dụng generate_text
        processed_questions = generate_text(generated_questions)
        
        print("Processed Questions:\n", processed_questions)

        save_as_pdf(processed_questions, output_file)

        print(f"Essay questions generated successfully and saved as {output_file}.")
    except FileNotFoundError as fnfe:
        print(f"Error: {fnfe}")
    except ValueError as ve:
        print(f"Value Error: {ve}")
    except Exception as e:
        print(f"An unexpected error occurred: {e}")

def process_file(file_path, num_questions, output_file):
    extracted_content = extract_text_from_file(file_path)
    full_content = "\n".join([doc.page_content for doc in extracted_content])

    generated_questions = generate_essay_questions(full_content, num_questions)
    save_as_pdf(generated_questions, output_file)

    return output_file


if __name__ == "__main__":
    main()
