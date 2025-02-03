import os
import logging
import uuid
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import FileResponse
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_community.document_loaders import PyPDFLoader

from docx import Document as DocxDocument
from fpdf import FPDF
from pptx import Presentation

# Configure logging
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")

# Load API Key from environment variables
os.environ["GOOGLE_API_KEY"] = "AIzaSyC437JOnEhmdcIqtuHSJAEkzaD6faXwy-I"

# Initialize AI model
llm = ChatGoogleGenerativeAI(
    model="gemini-1.5-pro",
    temperature=0.7,
    max_tokens=1024
)

# Prompt for generating lession examples
lession_prompt = PromptTemplate(
    input_variables=["content"],
    template=(
        "You are an experienced teacher. Below is a lession topic:\n"
        "{content}\n\n"
        "Create a structured document to help students understand the topic better:\n"
        "- Summarize the lession.\n"
        "- Provide at least two real-world examples illustrating key concepts.\n"
        "- Explain how this lession applies to real-life scenarios or work environments.\n"
        "Ensure clarity and well-structured formatting."
    )
)

# Function to read a .docx file
def load_docx(file_path):
    try:
        doc = DocxDocument(file_path)
        content = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
        return content if content.strip() else None
    except Exception as e:
        logging.error(f"Error reading .docx file: {e}")
        return None

# Function to read a .pdf file
def load_pdf(file_path):
    try:
        loader = PyPDFLoader(file_path)
        documents = loader.load()
        content = "\n".join(doc.page_content for doc in documents if doc.page_content.strip())
        return content if content.strip() else None
    except Exception as e:
        logging.error(f"Error reading PDF file: {e}")
        return None

# Function to read a .pptx file
def load_pptx(file_path):
    try:
        prs = Presentation(file_path)
        content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    content.append(shape.text)
        return "\n".join(content) if content else None
    except Exception as e:
        logging.error(f"Error reading PPTX file: {e}")
        return None

# Function to determine file type and extract content
def extract_text_from_file(file_path):
    if file_path.endswith(".docx"):
        return load_docx(file_path)
    elif file_path.endswith(".pdf"):
        return load_pdf(file_path)
    elif file_path.endswith(".pptx"):
        return load_pptx(file_path)
    else:
        raise ValueError("Unsupported file format. Only .docx, .pdf, and .pptx are supported.")

# Function to generate lession examples using AI
def generate_lession_content(content):
    try:
        lession_chain = LLMChain(llm=llm, prompt=lession_prompt)
        return lession_chain.run(content).strip()
    except Exception as e:
        logging.error(f"Error generating lession content: {e}")
        return None

# Function to save content as a PDF file
def save_as_pdf(output_text, file_name):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", style="B", size=16)
    pdf.set_text_color(0, 0, 128)
    pdf.cell(0, 10, "Connecting Lession Topics with Real-World Applications", ln=True, align="C")
    pdf.ln(10)

    pdf.set_text_color(0, 0, 0)
    pdf.set_font("Arial", style="", size=12)
    pdf.multi_cell(0, 10, output_text)

    pdf.output(file_name)
    logging.info(f"PDF saved at {file_name}")

# Function to process lession file
def process_lession_file(file_path, output_file):
    try:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        # Extract lession content from the file
        lession_content = extract_text_from_file(file_path)
        if not lession_content:
            raise ValueError("Extracted content is empty or invalid.")

        # Generate lession examples
        lession_output = generate_lession_content(lession_content)
        if not lession_output:
            raise ValueError("Failed to generate lession content.")

        logging.info("Lession content generated successfully!")

        # Save output to PDF
        save_as_pdf(lession_output, output_file)

        return output_file

    except Exception as e:
        logging.error(f"Error processing lession file: {e}")
        return None

