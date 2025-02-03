import os
import fitz  # PyMuPDF để đọc PDF
from pathlib import Path
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_google_genai import ChatGoogleGenerativeAI

# --- Cấu hình API Google Gemini ---
os.environ["GOOGLE_API_KEY"] = "AIzaSyC437JOnEhmdcIqtuHSJAEkzaD6faXwy-I"

# --- Khởi tạo mô hình Gemini ---
llm = ChatGoogleGenerativeAI(
    model="gemini-1.5-pro",
    temperature=0.3,
    max_tokens=4096
)



# --- Hàm đọc file PDF ---
def extract_text_from_pdf(file_path):
    """Trích xuất văn bản từ file PDF bằng PyMuPDF."""
    text = []
    try:
        with fitz.open(file_path) as pdf:
            for page_num, page in enumerate(pdf, start=1):
                page_text = page.get_text("text").strip()
                if page_text:
                    text.append(f"--- Page {page_num} ---\n{page_text}")
        return "\n\n".join(text) if text else "No text found in PDF."
    except Exception as e:
        raise RuntimeError(f"Error reading PDF: {e}")

# --- Hàm đọc file DOCX ---
def load_docx(file_path):
    """Trích xuất văn bản từ file DOCX."""
    try:
        doc = DocxDocument(file_path)
        return "\n".join([p.text.strip() for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        raise RuntimeError(f"Error reading DOCX: {e}")

# --- Hàm đọc file PPTX ---
def load_pptx(file_path):
    """Trích xuất văn bản từ file PPTX."""
    try:
        presentation = Presentation(file_path)
        content = []
        for slide_num, slide in enumerate(presentation.slides, start=1):
            print(f"Processing slide {slide_num}...")
            slide_content = []
            for shape in slide.shapes:
                print(f"  Shape type: {type(shape)}")
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        print(f"  Found text: {text[:50]}...")  # Xem trước văn bản
                        slide_content.append(text)
            if slide_content:
                content.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_content))
        return "\n\n".join(content) if content else "No text found in PPTX."
    except Exception as e:
        raise RuntimeError(f"Error reading PPTX: {e}")


# --- Định nghĩa các loại file hỗ trợ ---
FILE_LOADERS = {
    ".pdf": extract_text_from_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
}

def extract_text_from_file(file_path):
    """Lấy nội dung văn bản từ file."""
    file_extension = Path(file_path).suffix.lower()
    if file_extension not in FILE_LOADERS:
        raise ValueError(f"Unsupported file format: {file_extension}")
    return FILE_LOADERS[file_extension](file_path).strip()

# --- Gọi API Gemini để tạo mindmap ---
def generate_mindmap_with_gemini(content):
    """Gửi nội dung lên Gemini để tạo mindmap Markdown."""
    prompt = f"""
    Bạn là một chuyên gia trong việc tổ chức nội dung thành mindmap.
    Hãy phân tích văn bản sau và chuyển đổi nó thành Markdown mindmap có phân cấp rõ ràng.
    
    # Yêu cầu:
    - Sử dụng `#`, `##`, `###` để biểu diễn cấp độ của tiêu đề.
    - Dùng `-` hoặc `  -` để thể hiện các mục con.
    - Nếu có danh sách, giữ nguyên dạng danh sách.
    
    # Văn bản:
    {content}
    """
    
    response = llm.invoke(prompt)
    return response.content.strip() if response.content else "Error generating mindmap."

# --- Xử lý file input và tạo mindmap Markdown ---
def process_mindmap(input_file, output_md):
    """Trích xuất nội dung từ file và tạo mindmap Markdown sử dụng Gemini AI."""
    print(f"Processing mindmap for file: {input_file}")
    try:
        # Trích xuất nội dung từ file (PDF, DOCX, PPTX)
        extracted_content = extract_text_from_file(input_file)
        if not extracted_content:
            raise ValueError("No content extracted from the file.")

        print(f"Extracted content preview:\n{extracted_content[:500]}...")

        # Gửi lên Gemini để tạo Markdown
        markdown_content = generate_mindmap_with_gemini(extracted_content)
        if "Error" in markdown_content:
            raise ValueError("Failed to generate structured mindmap.")

        # Ghi vào file Markdown
        with open(output_md, "w", encoding="utf-8") as f:
            f.write(markdown_content)

        print(f"Mindmap saved as Markdown: {output_md}")
        return f"Success: Mindmap saved to {output_md}"

    except Exception as e:
        print(f"Error: {e}")
        return f"Error: {str(e)}"
