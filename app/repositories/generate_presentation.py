import os
import logging
import requests
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from docx import Document as DocxDocument
import fitz  # PyMuPDF
from transformers import MarianMTModel, MarianTokenizer
import re

# Logging configuration
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")

# Pexels API key
PEXELS_API_KEY = "K20qwM7iiTtdBMYHWBSeuv89nari4KvkI2DpGSlR0RZCQcRbqrJlk4gA"

# --- File Loaders ---
def load_docx(file_path):
    """Load text content from a DOCX file."""
    doc = DocxDocument(file_path)
    content = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
    if not content.strip():
        raise ValueError(f"File {file_path} is empty or invalid.")
    return content

def load_pdf(file_path):
    """Load text content from a PDF file."""
    doc = fitz.open(file_path)
    content = []
    for page in doc:
        content.append(page.get_text())
    doc.close()
    if not content:
        raise ValueError(f"File {file_path} is empty or invalid.")
    return "\n".join(content)

FILE_LOADERS = {
    ".docx": load_docx,
    ".pdf": load_pdf,
}

def extract_text_from_file(file_path):
    """Extract text from a supported file format."""
    file_extension = Path(file_path).suffix.lower()
    if file_extension not in FILE_LOADERS:
        raise ValueError(f"Unsupported file format: {file_extension}")
    loader = FILE_LOADERS[file_extension]
    return loader(file_path)

# --- Clean Text ---
def clean_text(text):
    """Clean input text to remove unnecessary characters like numbering."""
    text = re.sub(r"^\d+\.\s*", "", text)  # Loại bỏ số thứ tự
    text = re.sub(r"[^a-zA-Z0-9\s.,:()'-]", "", text)  # Loại bỏ ký tự đặc biệt không cần thiết
    return text.strip()

# --- Translate using Hugging Face ---
def translate_with_huggingface(text, src_lang="auto", tgt_lang="en"):
    """Translate text using Hugging Face MarianMT."""
    if src_lang == "auto":
        # Default to Vietnamese if language detection isn't implemented
        src_lang = "vi"
    model_name = f"Helsinki-NLP/opus-mt-{src_lang}-{tgt_lang}"
    tokenizer = MarianTokenizer.from_pretrained(model_name)
    model = MarianMTModel.from_pretrained(model_name)

    inputs = tokenizer([text], return_tensors="pt", padding=True, truncation=True)
    translated = model.generate(**inputs)
    return tokenizer.decode(translated[0], skip_special_tokens=True)

# --- Search Image using Pexels API ---
def search_image_pexels(query):
    """Search for an image using Pexels API."""
    url = f"https://api.pexels.com/v1/search?query={query}&per_page=1"
    headers = {"Authorization": PEXELS_API_KEY}
    try:
        response = requests.get(url, headers=headers)
        data = response.json()
        if "photos" in data and len(data["photos"]) > 0:
            return data["photos"][0]["src"]["medium"]  # Return the first image URL
        else:
            logging.warning(f"No images found for query: {query}")
            return None
    except Exception as e:
        logging.error(f"Error searching image for query '{query}' with Pexels API: {e}")
        return None

# --- Split Content by Sections ---
def split_into_sections(content):
    """Split content into logical sections based on headings."""
    sections = []
    current_section = []
    for line in content.splitlines():
        line = line.strip()
        if line and (line[0].isdigit() or line.isupper()):  # Detect headings
            if current_section:  # Save previous section
                sections.append("\n".join(current_section))
            current_section = [line]  # Start a new section
        else:
            current_section.append(line)
    if current_section:  # Save the last section
        sections.append("\n".join(current_section))
    return sections

# --- Process File to PowerPoint ---
def process_file_to_pptx(input_path: str, output_path: str):
    """Generate a PowerPoint presentation styled like the provided template."""
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Generated Presentation"
    title.text_frame.paragraphs[0].font.name = "Arial Black"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(33, 150, 243)  # Blue color

    subtitle.text = "Automotive Industry Style"
    subtitle.text_frame.paragraphs[0].font.name = "Arial"
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

    # Extract content
    content = extract_text_from_file(input_path)

    # Split content into sections
    sections = split_into_sections(content)

    # Create slides for each section
    for idx, section in enumerate(sections, 1):
        slide = prs.slides.add_slide(content_slide_layout)
        slide_title = slide.shapes.title
        slide_content = slide.placeholders[1]

        # First line of section as slide title
        lines = section.split("\n")
        slide_title.text = lines[0]
        slide_title.text_frame.paragraphs[0].font.name = "Arial Black"
        slide_title.text_frame.paragraphs[0].font.size = Pt(28)
        slide_title.text_frame.paragraphs[0].font.bold = True
        slide_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(33, 150, 243)

        # Remaining content as slide body
        slide_content.text = "\n".join(lines[1:])
        for paragraph in slide_content.text_frame.paragraphs:
            paragraph.font.name = "Arial"
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)

        # Search for an image for this slide
        translated_query = translate_with_huggingface(lines[0])
        logging.info(f"Translated query: {translated_query}")

        image_url = search_image_pexels(translated_query)
        if image_url:
            image_path = f"temp_image_{idx}.png"
            try:
                response = requests.get(image_url)
                with open(image_path, "wb") as img_file:
                    img_file.write(response.content)

                # Add image to the bottom-right corner of the slide
                slide.shapes.add_picture(image_path, Inches(5), Inches(3), width=Inches(4), height=Inches(3))
                os.remove(image_path)  # Clean up temporary file
            except Exception as e:
                logging.error(f"Error downloading or adding image for slide {idx}: {e}")

    prs.save(output_path)
    logging.info(f"Presentation successfully saved at {output_path}")
